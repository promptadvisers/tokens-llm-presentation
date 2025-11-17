from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Black and white minimalistic color palette
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY_LIGHT = RGBColor(240, 240, 240)
GRAY_MEDIUM = RGBColor(128, 128, 128)
GRAY_DARK = RGBColor(64, 64, 64)

def add_standard_slide(title_text, content_items):
    """Helper function to create standard content slides"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(15), Inches(1))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = BLACK
    
    # Minimalistic accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.color.rgb = BLACK
    
    # Content bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(15), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = GRAY_DARK
        p.space_before = Pt(12)
        p.level = 0
    
    return slide

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Main title
title_box = slide1.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Understanding Tokens"
title_p.font.size = Pt(66)
title_p.font.bold = True
title_p.font.color.rgb = BLACK
title_p.alignment = 1  # Center

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(5.2), Inches(14), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_p = subtitle_frame.paragraphs[0]
subtitle_p.text = "The Building Blocks of Large Language Models"
subtitle_p.font.size = Pt(28)
subtitle_p.font.color.rgb = GRAY_MEDIUM
subtitle_p.alignment = 1  # Center

# Minimalistic geometric element
rect = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7), Inches(6.8), Inches(2), Inches(0.1))
rect.fill.solid()
rect.fill.fore_color.rgb = BLACK
rect.line.color.rgb = BLACK

# Slide 2: What is a Token?
add_standard_slide(
    "What is a Token?",
    [
        "A token is the basic unit of text that a language model processes",
        "Tokens can be words, subwords, characters, or punctuation marks",
        "Models don't read text the way humans do—they process tokens",
        "Example: 'Hello world!' might be split into ['Hello', ' world', '!']"
    ]
)

# Slide 3: Why Tokens Matter
add_standard_slide(
    "Why Tokens Matter",
    [
        "Language models work with numbers, not text directly",
        "Text must be converted into tokens, then into numerical representations",
        "The tokenization method affects model performance and capabilities",
        "Token limits define how much text a model can process at once",
        "Understanding tokens helps optimize prompts and manage costs"
    ]
)

# Slide 4: Tokenization Process
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "The Tokenization Process"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = BLACK

# Accent line
line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(2), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = BLACK

# Process flow boxes
steps = [
    ("1. Raw Text", 2),
    ("2. Tokenize", 5),
    ("3. Token IDs", 8),
    ("4. Embeddings", 11)
]

for step, x_pos in steps:
    box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(3.5), Inches(2.5), Inches(1.5))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LIGHT
    box.line.color.rgb = BLACK
    box.line.width = Pt(2)
    
    text_frame = box.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = step
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = 1  # Center
    text_frame.vertical_anchor = 1  # Middle

# Arrows between boxes
arrow_positions = [(4.5, 7.5), (7.5, 10.5), (10.5, 13.5)]
for start_x, end_x in arrow_positions:
    arrow = slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(start_x), Inches(3.9), Inches(0.8), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLACK
    arrow.line.color.rgb = BLACK

# Example at bottom
example_box = slide4.shapes.add_textbox(Inches(2), Inches(6), Inches(12), Inches(2))
text_frame = example_box.text_frame
p = text_frame.paragraphs[0]
p.text = 'Example: "Hello world" → ["Hello", " world"] → [5158, 1917] → [vector embeddings]'
p.font.size = Pt(18)
p.font.color.rgb = GRAY_DARK
p.alignment = 1

# Slide 5: Types of Tokenization
add_standard_slide(
    "Types of Tokenization",
    [
        "Word-level: Each word becomes a token (simple but large vocabulary)",
        "Character-level: Each character is a token (flexible but long sequences)",
        "Subword: Balance between words and characters (most common)",
        "Byte-Pair Encoding (BPE): Merges frequent character pairs iteratively",
        "WordPiece & SentencePiece: Variations used by different models"
    ]
)

# Slide 6: Subword Tokenization Example
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Subword Tokenization Example"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = BLACK

# Accent line
line = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(2), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = BLACK

# Examples
examples = [
    ('Common word: "running"', '["running"]', 3),
    ('Uncommon word: "tokenization"', '["token", "ization"]', 4.5),
    ('Rare word: "antidisestablishmentarianism"', '["anti", "dis", "establish", "ment", "arian", "ism"]', 6)
]

y_pos = 2.8
for word, tokens, line_y in examples:
    # Word
    word_box = slide6.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(6), Inches(0.6))
    p = word_box.text_frame.paragraphs[0]
    p.text = word
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    
    # Arrow
    arrow = slide6.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(7.5), Inches(y_pos + 0.1), Inches(1), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GRAY_MEDIUM
    arrow.line.color.rgb = GRAY_MEDIUM
    
    # Tokens
    tokens_box = slide6.shapes.add_textbox(Inches(9), Inches(y_pos), Inches(6), Inches(0.6))
    p = tokens_box.text_frame.paragraphs[0]
    p.text = tokens
    p.font.size = Pt(20)
    p.font.color.rgb = GRAY_DARK
    p.font.name = "Courier New"
    
    y_pos += 1.4

# Key insight
insight_box = slide6.shapes.add_textbox(Inches(1), Inches(7), Inches(14), Inches(1.2))
text_frame = insight_box.text_frame
p = text_frame.paragraphs[0]
p.text = "Key Insight: Frequent words = fewer tokens, Rare words = more tokens"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = BLACK
p.alignment = 1

# Slide 7: Token Limits and Context Windows
add_standard_slide(
    "Token Limits & Context Windows",
    [
        "Every model has a maximum context window (measured in tokens)",
        "Context window includes both input (prompt) and output (response)",
        "Examples: GPT-3.5 (4K tokens), GPT-4 (8K-32K), Claude (200K)",
        "Exceeding limits requires truncation or summarization",
        "Longer contexts enable more complex reasoning and document analysis"
    ]
)

# Slide 8: Practical Implications
add_standard_slide(
    "Practical Implications",
    [
        "Cost: Many APIs charge per token (input + output)",
        "Speed: More tokens = longer processing time",
        "Context management: Must fit prompts within token limits",
        "Language differences: Some languages use more tokens than others",
        "Special characters and code often require more tokens than plain text"
    ]
)

# Slide 9: Optimizing Token Usage
add_standard_slide(
    "Optimizing Token Usage",
    [
        "Be concise: Remove unnecessary words from prompts",
        "Use clear structure: Well-organized text tokenizes more efficiently",
        "Choose the right model: Balance token limits with task requirements",
        "Monitor usage: Track token consumption for cost management",
        "Consider chunking: Break large documents into smaller segments",
        "Test tokenization: Use tokenizer tools to preview splits"
    ]
)

# Slide 10: Key Takeaways
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide10.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Title
title_box = slide10.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.paragraphs[0]
title_p.text = "Key Takeaways"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = BLACK

# Accent line
line = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.9), Inches(2), Inches(0.05))
line.fill.solid()
line.fill.fore_color.rgb = BLACK

# Key points in boxes
takeaways = [
    "Tokens are the fundamental units LLMs process",
    "Tokenization affects performance, cost, and capabilities",
    "Understanding tokens helps optimize AI interactions",
    "Different models use different tokenization strategies"
]

y_pos = 3
for takeaway in takeaways:
    box = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y_pos), Inches(12), Inches(0.9))
    box.fill.solid()
    box.fill.fore_color.rgb = GRAY_LIGHT
    box.line.color.rgb = BLACK
    box.line.width = Pt(1.5)
    
    text_frame = box.text_frame
    p = text_frame.paragraphs[0]
    p.text = takeaway
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.alignment = 1
    text_frame.vertical_anchor = 1
    
    y_pos += 1.2

# Save presentation
output_path = "/Users/marwankashef/Desktop/YouTube/Act III Demos/tokens_in_llms.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
