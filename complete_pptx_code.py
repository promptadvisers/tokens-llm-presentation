#!/usr/bin/env python3
"""
Complete PowerPoint Creation Script
This creates a professional 10-slide presentation about Language Models and Prompt Engineering
"""

# STEP 1: Install dependencies (run this in terminal first)
# pip install python-pptx pillow --break-system-packages

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Initialize presentation with widescreen format
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Define professional color palette
PRIMARY = RGBColor(30, 41, 59)      # Dark slate
ACCENT = RGBColor(99, 102, 241)     # Indigo
GRAY = RGBColor(148, 163, 184)      # Slate gray
WHITE = RGBColor(255, 255, 255)     # White

# SLIDE 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add diagonal accent shape
accent_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 
    Inches(-2), Inches(-2), 
    Inches(10), Inches(6)
)
accent_shape.fill.solid()
accent_shape.fill.fore_color.rgb = ACCENT
accent_shape.line.fill.background()
accent_shape.rotation = 15

# Add title
title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(2))
title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "Language Models & Prompt Engineering"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(10), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.clear()
p = subtitle_frame.add_paragraph()
p.text = "Understanding the Architecture and Art of AI Communication"
p.font.size = Pt(24)
p.font.color.rgb = GRAY

# SLIDE 2: What are Language Models?
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Add accent line
line = slide.shapes.add_connector(1, Inches(0), Inches(0.5), Inches(16), Inches(0.5))
line.line.width = Pt(3)
line.line.color.rgb = ACCENT

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
p = title_box.text_frame.add_paragraph()
p.text = "What are Language Models?"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Add content (narrower to make room for visual)
content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(8), Inches(5.5))
content = [
    "• AI systems trained on vast amounts of text data",
    "• Predict the next word in a sequence based on patterns",
    "• Built using neural network architectures",
    "• Learn statistical relationships between words",
    "• Can generate and understand human language"
]
for item in content:
    p = content_box.text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(16)

# Add neural network visualization
x_start = Inches(10)
y_start = Inches(3)

# Input layer (3 nodes)
for i in range(3):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x_start, y_start + Inches(i * 1.2),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT

# Hidden layer (4 nodes)
for i in range(4):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x_start + Inches(2), y_start - Inches(0.6) + Inches(i * 1.2),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = GRAY

# Output layer (2 nodes)
for i in range(2):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x_start + Inches(4), y_start + Inches(0.6) + Inches(i * 1.2),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY

# SLIDE 3: The Transformer Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
p = title_box.text_frame.add_paragraph()
p.text = "The Transformer Architecture"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

text_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
p = text_box.text_frame.add_paragraph()
p.text = "Modern language models use the Transformer architecture, which processes text through multiple layers of attention mechanisms. Each layer helps the model understand relationships between different parts of the input text."
p.font.size = Pt(18)
p.font.color.rgb = PRIMARY
p.word_wrap = True

# Transformer visual
x_pos = Inches(9.5)
y_pos = Inches(2.5)

components = [
    ("Input", RGBColor(248, 250, 252), ACCENT, PRIMARY, 0),
    ("Attention", ACCENT, ACCENT, WHITE, 1.2),
    ("Attention", ACCENT, ACCENT, WHITE, 2.1),
    ("Attention", ACCENT, ACCENT, WHITE, 3.0),
    ("Output", PRIMARY, PRIMARY, WHITE, 4.5)
]

for text, fill_color, line_color, text_color, y_offset in components:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x_pos, y_pos + Inches(y_offset),
        Inches(2), Inches(0.7 if text == "Attention" else 0.8)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.color.rgb = line_color
    tf = box.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12 if text == "Attention" else 14)
    p.font.color.rgb = text_color

# Helper function for standard slides
def add_standard_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Accent line
    line = slide.shapes.add_connector(1, Inches(0), Inches(0.5), Inches(16), Inches(0.5))
    line.line.width = Pt(3)
    line.line.color.rgb = ACCENT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
    p = title_box.text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(13), Inches(5.5))
    for point in bullet_points:
        p = content_box.text_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(16)
    
    return slide

# SLIDE 4: How LLMs Process Text
add_standard_slide(
    "How LLMs Process Text",
    [
        "1. Tokenization: Breaking text into smaller units",
        "2. Embedding: Converting tokens to numerical vectors",
        "3. Attention: Understanding relationships between tokens",
        "4. Layer Processing: Refining understanding through multiple layers",
        "5. Generation: Producing probability distributions for next tokens"
    ]
)

# SLIDE 5: Training Process
add_standard_slide(
    "Training Process",
    [
        "• Pre-training: Learning from massive text datasets",
        "• Self-supervised learning: Predicting masked or next words",
        "• Fine-tuning: Adapting to specific tasks or behaviors",
        "• Reinforcement Learning: Aligning with human preferences",
        "• Continuous improvement through feedback"
    ]
)

# SLIDE 6: Why Prompt Engineering Matters
add_standard_slide(
    "Why Prompt Engineering Matters",
    [
        "• LLMs are sensitive to input phrasing and structure",
        "• Different prompts activate different learned patterns",
        "• Context and instructions shape model behavior",
        "• Quality of output directly relates to prompt quality",
        "• Bridges the gap between human intent and AI understanding"
    ]
)

# SLIDE 7: Key Prompt Engineering Principles
add_standard_slide(
    "Key Prompt Engineering Principles",
    [
        "• Be Specific: Clear, detailed instructions yield better results",
        "• Provide Context: Background information improves relevance",
        "• Use Examples: Few-shot learning guides desired output format",
        "• Structure Matters: Organized prompts produce organized responses",
        "• Iterate and Refine: Test and improve prompts based on results"
    ]
)

# SLIDE 8: Comparison - Poor vs. Effective Prompts
slide = prs.slides.add_slide(prs.slide_layouts[6])

title_box = slide.shapes.add_textbox(Inches(1), Inches(0.8), Inches(14), Inches(1))
p = title_box.text_frame.add_paragraph()
p.text = "Poor vs. Effective Prompts"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = PRIMARY

# Left column header
left_header = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(2.3),
    Inches(6), Inches(0.8)
)
left_header.fill.solid()
left_header.fill.fore_color.rgb = GRAY
left_header.line.fill.background()
p = left_header.text_frame.add_paragraph()
p.text = "Poor Prompts"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = WHITE

# Left content
left_content = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(6), Inches(4))
left_items = [
    "• Vague instructions",
    "• Missing context",
    "• Ambiguous goals",
    "• No format specification",
    "• Single attempt"
]
for item in left_items:
    p = left_content.text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)

# Right column header
right_header = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(8.5), Inches(2.3),
    Inches(6), Inches(0.8)
)
right_header.fill.solid()
right_header.fill.fore_color.rgb = ACCENT
right_header.line.fill.background()
p = right_header.text_frame.add_paragraph()
p.text = "Effective Prompts"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = WHITE

# Right content
right_content = slide.shapes.add_textbox(Inches(8.5), Inches(3.5), Inches(6), Inches(4))
right_items = [
    "• Clear, specific directions",
    "• Relevant background info",
    "• Well-defined objectives",
    "• Output format examples",
    "• Iterative refinement"
]
for item in right_items:
    p = right_content.text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = PRIMARY
    p.space_after = Pt(12)

# SLIDE 9: Advanced Prompt Techniques
add_standard_slide(
    "Advanced Prompt Techniques",
    [
        "• Chain-of-Thought: Encouraging step-by-step reasoning",
        "• Role Playing: Defining personas for specialized responses",
        "• System Prompts: Setting behavioral guidelines",
        "• Temperature Control: Adjusting creativity vs. consistency",
        "• Prompt Chaining: Breaking complex tasks into steps"
    ]
)

# SLIDE 10: The Future of Language Models
add_standard_slide(
    "The Future of Language Models",
    [
        "• Multimodal capabilities: Text, images, audio, and video",
        "• Improved reasoning and mathematical abilities",
        "• Better alignment with human values and intent",
        "• More efficient architectures and training methods",
        "• Democratization of AI through better interfaces"
    ]
)

# SAVE THE PRESENTATION
output_path = "/mnt/user-data/outputs/language_models_prompt_engineering.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully!")
print(f"📍 Location: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
